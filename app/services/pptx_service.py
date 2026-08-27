from pathlib import Path

from pptx import Presentation as PptxPresentation

from app.domain import slides


def extract_text(pptx_path: Path | str) -> str:
    """Extrai o texto limpo de todos os slides do PPTx enviado pelo Cliente VR.

    Usa o mesmo contrato do `pdf_service` (`app.domain.slides`), para que o prompt do LLM
    e o aterramento não precisem saber de qual formato o texto veio.
    """
    presentation = PptxPresentation(str(pptx_path))
    blocos: list[str] = []

    for index, slide in enumerate(presentation.slides, start=1):
        blocks = [block for shape in slide.shapes for block in _shape_text(shape)]
        content = "\n".join(block for block in blocks if block)
        if content.strip():
            blocos.append(slides.bloco(index, content.strip()))

    return slides.montar(blocos)


def _shape_text(shape) -> list[str]:
    """Texto de uma forma: caixas de texto, tabelas e formas agrupadas."""
    # Formas agrupadas contêm outras formas — precisa descer recursivamente.
    if shape.shape_type is not None and getattr(shape, "shapes", None) is not None:
        return [block for inner in shape.shapes for block in _shape_text(inner)]

    if getattr(shape, "has_table", False):
        return [
            " | ".join(cell.text.strip() for cell in row.cells)
            for row in shape.table.rows
        ]

    if getattr(shape, "has_text_frame", False):
        return [shape.text_frame.text]

    return []


def is_valid_pptx(pptx_path: Path | str) -> bool:
    """Confere se o arquivo é realmente um PPTx legível antes de persistir a sessão."""
    try:
        PptxPresentation(str(pptx_path))
        return True
    except Exception:
        return False
